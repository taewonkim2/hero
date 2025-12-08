from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    # 프레젠테이션 객체 생성
    prs = Presentation()

    # 슬라이드 레이아웃 정의 (0: Title, 1: Title and Content, etc.)
    TITLE_SLIDE_LAYOUT = prs.slide_layouts[0]
    BULLET_SLIDE_LAYOUT = prs.slide_layouts[1]
    SECTION_HEADER_LAYOUT = prs.slide_layouts[2]

    # --- Slide 1: 표지 ---
    slide = prs.slides.add_slide(TITLE_SLIDE_LAYOUT)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "AI 기반 High-Fidelity\nUI/UX 컨설팅 제안"
    subtitle.text = "상상하는 PPT가 아닌,\n동작하는 경험(HTML)을 납품합니다"

    # --- Slide 2: 목차 ---
    slide = prs.slides.add_slide(BULLET_SLIDE_LAYOUT)
    shapes = slide.shapes
    shapes.title.text = "목차"
    
    tf = shapes.placeholders[1].text_frame
    tf.text = "1. 제안 배경 (Executive Summary)"
    
    p = tf.add_paragraph()
    p.text = "2. 기존 방식 vs 제안 방식 비교"
    p = tf.add_paragraph()
    p.text = "3. 프로젝트 수행 방법론"
    p = tf.add_paragraph()
    p.text = "4. 기대 효과"
    p = tf.add_paragraph()
    p.text = "5. 결론"

    # --- Slide 3: 제안 배경 ---
    slide = prs.slides.add_slide(BULLET_SLIDE_LAYOUT)
    shapes = slide.shapes
    shapes.title.text = "1. 제안 배경: 왜 앱 프로젝트는 결과물이 다를까?"
    
    tf = shapes.placeholders[1].text_frame
    tf.text = "Current Problem (문제점)"
    
    p = tf.add_paragraph()
    p.text = "죽은 문서(PPT)의 한계: 실제 동작과 느낌을 알 수 없음"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "주객전도된 기획: 개발 편의성 위주의 타협적인 UI 양산"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "검증 불가능: 텍스트 설명만으로는 UX 검증 불가"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Our Solution (해결책)"
    
    p = tf.add_paragraph()
    p.text = "Living Output: 동작하는 HTML/Web 시뮬레이터 제공"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Zero-Gap: 기획과 개발 사이의 괴리를 0으로 만듦"
    p.level = 1

    # --- Slide 4: As-Is vs To-Be ---
    slide = prs.slides.add_slide(BULLET_SLIDE_LAYOUT)
    shapes = slide.shapes
    shapes.title.text = "2. 기존 방식(As-Is) vs 제안 방식(To-Be)"
    
    tf = shapes.placeholders[1].text_frame
    tf.text = "기존 UI/UX 컨설팅"
    p = tf.add_paragraph()
    p.text = "산출물: 정적인 PPT 화면 설계서"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "검증: 눈으로 읽고 머리로 상상"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "혁신성: 기존 레퍼런스 짜깁기"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "AI 인터랙티브 컨설팅 (New!)"
    p = tf.add_paragraph()
    p.text = "산출물: 동작하는 HTML/Web 시뮬레이터"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "검증: 직접 스마트폰으로 눌러보고 체감"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "혁신성: AI와 함께 무한한 UX 실험 및 최적화"
    p.level = 1

    # --- Slide 5: 수행 방법론 ---
    slide = prs.slides.add_slide(BULLET_SLIDE_LAYOUT)
    shapes = slide.shapes
    shapes.title.text = "3. 프로젝트 수행 방법론"
    
    tf = shapes.placeholders[1].text_frame
    
    tf.text = "Phase 1: AI 기반 급속 프로토타이핑 (Rapid Prototyping)"
    p = tf.add_paragraph()
    p.text = "와이어프레임 생략, 자연어로 즉시 HTML 생성 (Vibe Coding)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Phase 2: 인터랙티브 UX 검증 (Interactive Validation)"
    p = tf.add_paragraph()
    p.text = "고객사 담당자 폰으로 실시간 테스트 및 즉각적 피드백 반영"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Phase 3: 개발 최적화 가이드 (Dev Optimization)"
    p = tf.add_paragraph()
    p.text = "확정된 HTML 코드를 분석하여 역공학적 개발 가이드 제공"
    p.level = 1

    # --- Slide 6: 기대 효과 ---
    slide = prs.slides.add_slide(BULLET_SLIDE_LAYOUT)
    shapes = slide.shapes
    shapes.title.text = "4. 기대 효과"
    
    tf = shapes.placeholders[1].text_frame
    
    tf.text = "고객 경험(UX) 혁신"
    p = tf.add_paragraph()
    p.text = "실제 동작 검증으로 사용성 실패 확률 Zero 도전"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "프로젝트 기간 단축 및 비용 절감"
    p = tf.add_paragraph()
    p.text = "기획 수정 재작업 방지, 전체 일정 20~30% 단축"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "명확한 의사소통"
    p = tf.add_paragraph()
    p.text = "모호한 텍스트 대신 실제 화면으로 소통"
    p.level = 1

    # --- Slide 7: 결론 ---
    slide = prs.slides.add_slide(SECTION_HEADER_LAYOUT)
    shapes = slide.shapes
    shapes.title.text = "보지 않으면 믿을 수 없고,\n만지지 않으면 느낄 수 없습니다."
    shapes.placeholders[1].text = "PPT 100장보다 강력한\n'동작하는 프로토타입 1개'를 제안합니다."

    # 파일 저장
    file_name = "AI_Interactive_Consulting_Proposal.pptx"
    prs.save(file_name)
    print(f"'{file_name}' 파일이 성공적으로 생성되었습니다!")

if __name__ == "__main__":
    create_presentation()